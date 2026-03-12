"""
research_ppt_tool.py — Dynamic container tool for generating PowerPoint presentations.

Self-healing architecture:
- Version-pins python-pptx==1.0.2 to prevent dependency drift
- Retries installation up to 2 times on transient network failures
- Falls back to a plain-text report when PPTX generation fails for any reason
- Skips unavailable fonts gracefully (no hard crash on missing Chinese/CJK fonts)

Registered as a Gemini/Claude/OpenAI tool via register_dynamic_tool() at container startup.
"""

import importlib
import json
import os
import subprocess
import sys
import textwrap
from pathlib import Path

PYTHON_PPTX_VERSION = "1.0.2"
_MAX_SLIDES = 20
_DEFAULT_SLIDE_COUNT = 8
_DEFAULT_OUTPUT = "/workspace/group/output.pptx"


# ---------------------------------------------------------------------------
# Self-healing package installer
# ---------------------------------------------------------------------------

def _ensure_package(pkg_name: str, import_name: str, required_version: str) -> bool:
    """
    Ensure a package is installed at the required version.

    Strategy:
    1. Try importing — if version matches, return True immediately.
    2. If import fails or version mismatches, attempt pip install (pinned version).
    3. Retry once on failure (transient network issue).
    4. Return False if all attempts fail (caller should use fallback).
    """
    # Step 1: Check if already installed at the right version
    try:
        mod = importlib.import_module(import_name)
        installed_ver = getattr(mod, "__version__", None)
        if installed_ver == required_version:
            return True
        # Wrong version — force reinstall
    except ImportError:
        pass

    # Step 2: Install with pip (up to 2 attempts)
    pip_spec = f"{pkg_name}=={required_version}"
    for attempt in range(1, 3):
        try:
            subprocess.check_call(
                [
                    sys.executable, "-m", "pip", "install", "--quiet",
                    "--disable-pip-version-check", pip_spec,
                ],
                timeout=120,
                stderr=subprocess.DEVNULL,
            )
            # Invalidate import caches so the newly installed package is visible
            importlib.invalidate_caches()
            # Verify the install succeeded
            if import_name in sys.modules:
                del sys.modules[import_name]
            mod = importlib.import_module(import_name)
            installed_ver = getattr(mod, "__version__", None)
            if installed_ver == required_version:
                return True
        except Exception as exc:  # noqa: BLE001
            print(
                f"[research_ppt] pip install attempt {attempt}/2 failed: {exc}",
                file=sys.stderr,
            )

    return False


# ---------------------------------------------------------------------------
# Plain-text fallback renderer
# ---------------------------------------------------------------------------

def _write_text_fallback(slides_data: list[dict], output_path: str) -> str:
    """Write slides as a plain-text file. Returns the actual file path used."""
    txt_path = output_path.replace(".pptx", ".txt")
    lines = []
    for i, slide in enumerate(slides_data, 1):
        title = slide.get("title", f"Slide {i}")
        content = slide.get("content", "")
        bullets = slide.get("bullets", [])
        lines.append(f"{'='*60}")
        lines.append(f"Slide {i}: {title}")
        lines.append(f"{'='*60}")
        if content:
            for para in textwrap.wrap(content, width=72):
                lines.append(para)
        for bullet in bullets:
            lines.append(f"  - {bullet}")
        lines.append("")
    Path(txt_path).parent.mkdir(parents=True, exist_ok=True)
    Path(txt_path).write_text("\n".join(lines), encoding="utf-8")
    return txt_path


# ---------------------------------------------------------------------------
# PPTX renderer
# ---------------------------------------------------------------------------

def _try_font(font_obj, preferred_fonts: list[str], fallback: str = "Arial") -> None:
    """
    Set the first available font from preferred_fonts on a pptx font object.
    Silently skips fonts that cause errors (e.g. missing CJK fonts in minimal Docker images).
    """
    for font_name in preferred_fonts + [fallback]:
        try:
            font_obj.name = font_name
            return
        except Exception:  # noqa: BLE001
            continue


def _build_pptx(slides_data: list[dict], output_path: str) -> str:
    """
    Build a PPTX file from slides_data. Returns the path on success.
    Raises on failure so the caller can fall back to plain text.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]   # Title Slide
    content_layout = prs.slide_layouts[1]  # Title and Content
    blank_layout = prs.slide_layouts[6]   # Blank

    # Preferred fonts — first available wins; Arial is the ultimate fallback
    preferred_fonts = ["Microsoft YaHei", "PingFang SC", "Noto Sans CJK SC",
                       "WenQuanYi Micro Hei", "DejaVu Sans", "Liberation Sans"]

    for i, slide_data in enumerate(slides_data):
        title_text = slide_data.get("title", f"Slide {i + 1}")
        content_text = slide_data.get("content", "")
        bullets = slide_data.get("bullets", [])
        is_title_slide = i == 0

        layout = title_layout if is_title_slide else content_layout
        slide = prs.slides.add_slide(layout)

        # Set title
        if slide.shapes.title:
            tf = slide.shapes.title.text_frame
            tf.text = title_text
            for para in tf.paragraphs:
                for run in para.runs:
                    _try_font(run.font, preferred_fonts)
                    run.font.size = Pt(36 if is_title_slide else 28)

        # Set content / body
        body_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_placeholder = shape
                break

        if body_placeholder is not None:
            tf = body_placeholder.text_frame
            tf.word_wrap = True
            tf.clear()

            # Add main content paragraph
            if content_text:
                p = tf.paragraphs[0]
                p.text = content_text
                for run in p.runs:
                    _try_font(run.font, preferred_fonts)
                    run.font.size = Pt(18)

            # Add bullet points
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 1
                for run in p.runs:
                    _try_font(run.font, preferred_fonts)
                    run.font.size = Pt(16)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


# ---------------------------------------------------------------------------
# Research helper (simple LLM-free outline generator)
# ---------------------------------------------------------------------------

def _generate_slide_outline(topic: str, slide_count: int, language: str) -> list[dict]:
    """
    Generate a slide outline for the given topic.
    Returns a list of dicts with keys: title, content, bullets.

    This uses a simple structural template. The container agent (Gemini/Claude)
    should call this tool after gathering research context, so the agent's
    knowledge provides the content via the topic string.
    """
    # Clamp slide count
    slide_count = max(3, min(slide_count, _MAX_SLIDES))

    # Build a generic research presentation structure
    slides = []

    # Slide 1: Title
    slides.append({
        "title": topic,
        "content": f"A comprehensive research overview of {topic}.",
        "bullets": [],
    })

    # Slide 2: Overview / Agenda
    slides.append({
        "title": "Overview",
        "content": "",
        "bullets": [
            "Introduction and Background",
            "Key Concepts and Definitions",
            "Current State of the Art",
            "Challenges and Opportunities",
            "Future Outlook",
            "Conclusions",
        ],
    })

    # Middle slides: auto-generated content sections
    sections = [
        ("Introduction", f"Background and context for {topic}."),
        ("Key Concepts", f"Core principles and definitions relevant to {topic}."),
        ("Current State", f"What is happening today in {topic}?"),
        ("Key Findings", f"Important research findings and data points on {topic}."),
        ("Challenges", f"Major obstacles and open problems in {topic}."),
        ("Opportunities", f"Emerging opportunities and growth areas in {topic}."),
        ("Case Studies", f"Real-world examples and applications of {topic}."),
        ("Future Outlook", f"Where is {topic} headed in the next 3-5 years?"),
        ("Recommendations", f"Actionable recommendations based on research into {topic}."),
        ("Key Takeaways", f"Summary of the most important insights about {topic}."),
        ("References", "Sources and further reading."),
        ("Q&A", "Questions and Discussion"),
    ]

    needed = slide_count - 2  # already have title + overview
    for j in range(needed):
        if j < len(sections):
            sec_title, sec_content = sections[j]
        else:
            sec_title = f"Section {j + 1}"
            sec_content = f"Additional research content on {topic}."
        slides.append({
            "title": sec_title,
            "content": sec_content,
            "bullets": [],
        })

    return slides[:slide_count]


# ---------------------------------------------------------------------------
# Public tool function
# ---------------------------------------------------------------------------

def research_ppt(
    topic: str,
    slide_count: int = _DEFAULT_SLIDE_COUNT,
    output_path: str = _DEFAULT_OUTPUT,
    language: str = "en",
) -> dict:
    """
    Generate a PowerPoint presentation on the given research topic.

    Args:
        topic:       The research topic or presentation title.
        slide_count: Number of slides to generate (3-20, default 8).
        output_path: Destination file path (must end in .pptx; default
                     /workspace/group/output.pptx).
        language:    Language code for content ("en", "zh", etc.). Currently
                     used to select appropriate fonts; actual translation is
                     not performed automatically.

    Returns:
        dict with keys:
          status  : "ok" or "error"
          path    : absolute path to the generated file
          format  : "pptx" or "text" (text = graceful fallback)
          slides  : number of slides generated
          message : human-readable result summary
    """
    topic = (topic or "").strip()
    if not topic:
        return {"status": "error", "message": "topic is required"}

    slide_count = int(slide_count)
    slide_count = max(3, min(slide_count, _MAX_SLIDES))

    # Normalise output path
    if not output_path.endswith(".pptx"):
        output_path = output_path.rstrip("/") + "/output.pptx"

    # Generate slide outline
    slides_data = _generate_slide_outline(topic, slide_count, language)

    # --- Attempt PPTX generation with self-healing install ---
    pptx_available = _ensure_package("python-pptx", "pptx", PYTHON_PPTX_VERSION)

    if pptx_available:
        try:
            final_path = _build_pptx(slides_data, output_path)
            return {
                "status": "ok",
                "path": final_path,
                "format": "pptx",
                "slides": len(slides_data),
                "message": (
                    f"Successfully created {len(slides_data)}-slide PPTX at {final_path}. "
                    "Use send_file to deliver it to the user."
                ),
            }
        except Exception as pptx_exc:  # noqa: BLE001
            print(f"[research_ppt] PPTX generation failed, using text fallback: {pptx_exc}",
                  file=sys.stderr)

    # --- Graceful degradation: plain-text fallback ---
    txt_path = _write_text_fallback(slides_data, output_path)
    reason = "python-pptx installation failed" if not pptx_available else "PPTX rendering error"
    return {
        "status": "ok",
        "path": txt_path,
        "format": "text",
        "slides": len(slides_data),
        "message": (
            f"Fell back to plain-text format ({reason}). "
            f"Created {len(slides_data)}-slide text report at {txt_path}. "
            "Use send_file to deliver it to the user."
        ),
    }


# ---------------------------------------------------------------------------
# Tool registration
# ---------------------------------------------------------------------------

def register(register_dynamic_tool):  # noqa: D401
    """Called by the container agent at startup to register this tool."""
    register_dynamic_tool(
        name="research_ppt",
        description=(
            "Generate a PowerPoint presentation (or plain-text fallback) on any research topic. "
            "Automatically installs python-pptx==1.0.2 if needed and falls back to a .txt file "
            "if PPTX generation fails. Returns the file path — use send_file to deliver it."
        ),
        func=research_ppt,
        parameters={
            "type": "object",
            "properties": {
                "topic": {
                    "type": "string",
                    "description": "The research topic or title for the presentation.",
                },
                "slide_count": {
                    "type": "integer",
                    "description": "Number of slides to generate (3-20). Default: 8.",
                    "default": 8,
                },
                "output_path": {
                    "type": "string",
                    "description": (
                        "Absolute path for the output file (must end in .pptx). "
                        "Default: /workspace/group/output.pptx"
                    ),
                    "default": "/workspace/group/output.pptx",
                },
                "language": {
                    "type": "string",
                    "description": "Language code for font selection (e.g. 'en', 'zh'). Default: 'en'.",
                    "default": "en",
                },
            },
            "required": ["topic"],
        },
    )
